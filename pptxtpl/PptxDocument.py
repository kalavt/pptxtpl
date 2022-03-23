from jinja2 import Environment, exceptions
from pptx import Presentation


class PptxDocument:
    def __init__(self, input_path, data=None, jinja_env=None):
        self.ppt = Presentation(input_path)
        self._init(data, jinja_env)

    def _init(self, data=None, jinja_env=None):
        if data:
            self.data = data
        if jinja_env:
            self.jinja_env = jinja_env
        else:
            self.jinja_env = Environment()

    def render(self, data=None, jinja_env=None):
        self._init(data, jinja_env)
        for slide in self.ppt.slides:
            self._render_slide(slide)

    def save(self, file):
        self.ppt.save(file)

    def _render_slide(self, slide):
        for shape in slide.shapes:
            self._render_shape(shape)

    def _render_shape(self, shape):
        if shape.has_text_frame:
            self._render_text_frame(shape.text_frame)
        if shape.has_table:
            self._render_table(shape.table)

    def _render_table(self, table):
        for row in table.rows:
            for cell in row.cells:
                self._render_text_frame(cell.text_frame)

    def _render_text_frame(self, text_frame):
        for para in text_frame.paragraphs:
            if "{" in para.text:
                print(para.text)
                expand = self._render_text(para.text)
                self._replace_paragraph_text(para, expand)

    def _replace_paragraph_text(self, paragraph, new_text):
        if len(paragraph.runs):
            p = (
                paragraph._p
            )  # the lxml element containing the `<a:p>` paragraph element
            # remove all but the first run
            for idx, run in enumerate(paragraph.runs):
                if idx == 0:
                    continue
                p.remove(run._r)
            paragraph.runs[0].text = new_text

    def _fix_qoute(self, text):
        dic = {"“": '"', "”": '"', "‘": "'", "’": "'"}
        for i, j in dic.items():
            text = text.replace(i, j)
        return text

    def _render_text(self, text):
        template = self.jinja_env.from_string(self._fix_qoute(text))
        try:
            result = template.render(self.data)
        except exceptions.UndefinedError as error:
            print(str(error))
        except exceptions.TemplateError as error:
            print(str(error))
        else:
            return result
